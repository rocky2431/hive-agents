from __future__ import annotations

import logging
import os
import shutil
import subprocess

logger = logging.getLogger(__name__)
from pathlib import Path

from app.config import get_settings
from app.skills import SkillRegistry, WorkspaceSkillLoader
from app.tools.packs import iter_tool_packs

WORKSPACE_ROOT = Path(get_settings().AGENT_DATA_DIR)


def _list_files(ws: Path, rel_path: str, tenant_id: str | None = None) -> str:
    if rel_path and rel_path.startswith("enterprise_info"):
        if tenant_id:
            enterprise_root = (WORKSPACE_ROOT / f"enterprise_info_{tenant_id}").resolve()
        else:
            enterprise_root = (WORKSPACE_ROOT / "enterprise_info").resolve()
        sub = rel_path[len("enterprise_info"):].lstrip("/")
        target = (enterprise_root / sub).resolve() if sub else enterprise_root
        if not str(target).startswith(str(enterprise_root)):
            return "Access denied for this path"
    else:
        target = (ws / rel_path) if rel_path else ws
        target = target.resolve()
        if not str(target).startswith(str(ws.resolve())):
            return "Access denied for this path"

    if not target.exists():
        return f"Directory not found: {rel_path or '/'}"

    items = []
    if not rel_path:
        if tenant_id:
            enterprise_dir = WORKSPACE_ROOT / f"enterprise_info_{tenant_id}"
        else:
            enterprise_dir = WORKSPACE_ROOT / "enterprise_info"
        if enterprise_dir.exists():
            items.append("  📁 enterprise_info/ (shared company info)")

    dir_count = 0
    file_count = 0
    for p in sorted(target.iterdir()):
        if p.name.startswith("."):
            continue
        if p.is_dir():
            dir_count += 1
            child_count = len([c for c in p.iterdir() if not c.name.startswith(".")])
            items.append(f"  📁 {p.name}/ ({child_count} items)")
        elif p.is_file():
            file_count += 1
            size_bytes = p.stat().st_size
            size_str = f"{size_bytes}B" if size_bytes < 1024 else f"{size_bytes/1024:.1f}KB"
            items.append(f"  📄 {p.name} ({size_str})")

    if not items:
        return f"📂 {rel_path or 'root'}: Empty directory (0 files, 0 folders)"

    header = f"📂 {rel_path or 'root'}: {dir_count} folder(s), {file_count} file(s)\n"
    return header + "\n".join(items)


def _read_file(ws: Path, rel_path: str, tenant_id: str | None = None) -> str:
    if rel_path and rel_path.startswith("enterprise_info"):
        if tenant_id:
            enterprise_root = (WORKSPACE_ROOT / f"enterprise_info_{tenant_id}").resolve()
        else:
            enterprise_root = (WORKSPACE_ROOT / "enterprise_info").resolve()
        sub = rel_path[len("enterprise_info"):].lstrip("/")
        file_path = (enterprise_root / sub).resolve() if sub else enterprise_root
        if not str(file_path).startswith(str(enterprise_root)):
            return "Access denied for this path"
    else:
        file_path = (ws / rel_path).resolve()
        if not str(file_path).startswith(str(ws.resolve())):
            return "Access denied for this path"

    if not file_path.exists():
        return f"File not found: {rel_path}"

    try:
        content = file_path.read_text(encoding="utf-8", errors="replace")
        if len(content) > 16000:
            content = content[:16000] + f"\n\n...[truncated, {len(content)} chars total]"
        return content
    except Exception as e:
        return f"Read failed: {e}"


def _load_skill(ws: Path, skill_name: str) -> str:
    requested = (skill_name or "").strip()
    if not requested:
        return "❌ Skill name cannot be empty"

    skills_dir = (ws / "skills").resolve()
    if not skills_dir.exists():
        return "Skill not found: skills directory does not exist"

    def _read_skill_file(path: Path) -> str:
        if not str(path).startswith(str(skills_dir)):
            return "Access denied for this skill path"
        rel_path = path.relative_to(ws).as_posix()
        return _read_file(ws, rel_path)

    requested_path = requested
    if requested_path.startswith("skills/"):
        requested_path = requested_path[len("skills/"):]
    explicit_path = (skills_dir / requested_path).resolve()
    if explicit_path.is_file():
        return _read_skill_file(explicit_path)

    registry = _build_skill_registry(ws)
    try:
        return registry.load_body(requested)
    except KeyError:
        return f"Skill not found: {skill_name}"


def _build_skill_registry(ws: Path) -> SkillRegistry:
    loader = WorkspaceSkillLoader()
    registry = SkillRegistry()
    registry.register_many(loader.load_from_workspace(ws))
    return registry


async def _read_document(ws: Path, rel_path: str, max_chars: int = 8000, tenant_id: str | None = None) -> str:
    if rel_path and rel_path.startswith("enterprise_info"):
        if tenant_id:
            enterprise_root = (WORKSPACE_ROOT / f"enterprise_info_{tenant_id}").resolve()
        else:
            enterprise_root = (WORKSPACE_ROOT / "enterprise_info").resolve()
        sub = rel_path[len("enterprise_info"):].lstrip("/")
        file_path = (enterprise_root / sub).resolve() if sub else enterprise_root
        if not str(file_path).startswith(str(enterprise_root)):
            return "Access denied for this path"
    else:
        file_path = (ws / rel_path).resolve()
        if not str(file_path).startswith(str(ws.resolve())):
            return "Access denied for this path"

    if not file_path.exists():
        return f"File not found: {rel_path}"

    ext = file_path.suffix.lower()
    try:
        if ext == ".pdf":
            import pdfplumber
            text_parts = []
            with pdfplumber.open(str(file_path)) as pdf:
                for i, page in enumerate(pdf.pages[:50]):
                    page_text = page.extract_text() or ""
                    if page_text:
                        text_parts.append(f"--- Page {i+1} ---\n{page_text}")
            content = "\n\n".join(text_parts) if text_parts else "(PDF is empty or text extraction failed)"
        elif ext == ".docx":
            from docx import Document
            from docx.oxml.ns import qn
            doc = Document(str(file_path))
            lines: list[str] = []

            def _extract_table(table) -> str:
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    deduped = [cells[0]] + [c for i, c in enumerate(cells[1:]) if c != cells[i]]
                    row_str = " | ".join(c for c in deduped if c)
                    if row_str:
                        rows.append(row_str)
                return "\n".join(rows)

            for para in doc.paragraphs:
                t = para.text.strip()
                if t:
                    lines.append(t)
            for table in doc.tables:
                t = _extract_table(table)
                if t:
                    lines.append(t)
            for shape in doc.element.body.iter(qn("w:txbxContent")):
                for child in shape.iter(qn("w:t")):
                    if child.text and child.text.strip():
                        lines.append(child.text.strip())
            for section in doc.sections:
                for hf in [section.header, section.footer]:
                    if hf and hf.is_linked_to_previous is False:
                        for para in hf.paragraphs:
                            t = para.text.strip()
                            if t:
                                lines.append(t)
            content = "\n".join(lines) if lines else "(Document is empty or uses unsupported formatting)"
        elif ext == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(str(file_path), read_only=True, data_only=True)
            sheets = []
            for ws_name in wb.sheetnames[:10]:
                sheet = wb[ws_name]
                rows = []
                for row in sheet.iter_rows(max_row=200, values_only=True):
                    row_str = "\t".join(str(c) if c is not None else "" for c in row)
                    if row_str.strip():
                        rows.append(row_str)
                if rows:
                    sheets.append(f"=== Sheet: {ws_name} ===\n" + "\n".join(rows))
            wb.close()
            content = "\n\n".join(sheets) if sheets else "(Excel is empty)"
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(file_path))
            slides = []
            for i, slide in enumerate(prs.slides[:50]):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text)
                if texts:
                    slides.append(f"--- Slide {i+1} ---\n" + "\n".join(texts))
            content = "\n\n".join(slides) if slides else "(PPT is empty)"
        elif ext in (".txt", ".md", ".json", ".csv", ".log"):
            content = file_path.read_text(encoding="utf-8", errors="replace")
        else:
            return f"Unsupported file format: {ext}. Supported: PDF, DOCX, XLSX, PPTX, TXT, MD, CSV"

        if len(content) > max_chars:
            content = content[:max_chars] + f"\n\n...[truncated, {len(content)} chars total]"
        return content
    except ImportError as e:
        return f"Missing dependency: {e}. Install: pip install pdfplumber python-docx openpyxl python-pptx"
    except Exception as e:
        return f"Document read failed: {str(e)[:200]}"


_WRITE_PROTECTED = {
    "tasks.json": "tasks.json is read-only. Use manage_tasks tool to manage tasks.",
}

# soul.md is append-only: heartbeat can add evolution notes but not overwrite identity
_APPEND_ONLY = {"soul.md"}


def _write_file(ws: Path, rel_path: str, content: str) -> str:
    if not rel_path or not rel_path.strip("/"):
        return "❌ Missing file path. Usage: write_file(path='workspace/report.md', content='...')"

    _blocked = _WRITE_PROTECTED.get(rel_path.strip("/"))
    if _blocked:
        return _blocked

    # soul.md is append-only: new content is appended under an evolution section
    _APPEND_ONLY_MAX_CHARS = 16000
    if rel_path.strip("/") in _APPEND_ONLY:
        target = ws / rel_path.strip("/")
        if target.exists():
            existing = target.read_text(encoding="utf-8", errors="replace")
            if content.strip() in existing:
                return f"✅ {rel_path} already contains this content."
            # Enforce size cap — trim oldest evolution notes if exceeding limit
            if len(existing) + len(content) > _APPEND_ONLY_MAX_CHARS:
                separator = "\n\n---\n## Evolution Notes (heartbeat-appended)\n\n"
                if separator.rstrip() in existing:
                    identity, _, evo_notes = existing.partition(separator.rstrip())
                    # Keep identity + trim evolution notes from the top
                    trimmed_notes = evo_notes[len(content):]  # drop oldest chars equal to new content size
                    existing = identity + separator.rstrip() + trimmed_notes
            separator = "\n\n---\n## Evolution Notes (heartbeat-appended)\n\n"
            if separator.rstrip() in existing:
                target.write_text(existing.rstrip() + "\n\n" + content.strip() + "\n", encoding="utf-8")
            else:
                target.write_text(existing.rstrip() + separator + content.strip() + "\n", encoding="utf-8")
            return f"✅ Appended evolution notes to {rel_path} (identity section preserved)."
        # If file doesn't exist yet, fall through to normal write

    file_path = (ws / rel_path).resolve()
    if not str(file_path).startswith(str(ws.resolve())):
        return "Access denied for this path"

    try:
        file_path.parent.mkdir(parents=True, exist_ok=True)
        file_path.write_text(content, encoding="utf-8")
        return f"✅ Written to {rel_path} ({len(content)} chars)"
    except Exception as e:
        return f"Write failed: {e}"


def _edit_file(ws: Path, rel_path: str, old_text: str, new_text: str, replace_all: bool = False) -> str:
    file_path = (ws / rel_path).resolve()
    if not str(file_path).startswith(str(ws.resolve())):
        return "Access denied for this path"
    if not file_path.exists():
        return f"File not found: {rel_path}"

    try:
        original = file_path.read_text(encoding="utf-8", errors="replace")
        occurrences = original.count(old_text)
        if occurrences == 0:
            return f"❌ Could not find the target text in {rel_path}"
        if not replace_all and occurrences != 1:
            return (
                f"❌ Found {occurrences} matches in {rel_path}. "
                "Refine old_text or set replace_all=true."
            )
        updated = original.replace(old_text, new_text, -1 if replace_all else 1)
        file_path.write_text(updated, encoding="utf-8")
        replaced = occurrences if replace_all else 1
        return f"✅ Updated {rel_path} ({replaced} replacement{'s' if replaced != 1 else ''})"
    except Exception as e:
        return f"Edit failed: {e}"


def _glob_search(ws: Path, pattern: str, root: str = "") -> str:
    search_root = (ws / root).resolve() if root else ws.resolve()
    if not str(search_root).startswith(str(ws.resolve())):
        return "Access denied for this path"
    if not search_root.exists():
        return f"Directory not found: {root or '/'}"

    matches: list[str] = []
    try:
        for path in sorted(search_root.glob(pattern)):
            resolved = path.resolve()
            if not str(resolved).startswith(str(ws.resolve())):
                continue
            matches.append(resolved.relative_to(ws).as_posix())
            if len(matches) >= 100:
                break
    except Exception as e:
        return f"Glob search failed: {e}"

    if not matches:
        return f"🔎 No files matched pattern '{pattern}'"
    lines = [f"🔎 Glob results for '{pattern}' ({len(matches)} match(es)):"]
    lines.extend(f"- {match}" for match in matches)
    return "\n".join(lines)


def _grep_search(ws: Path, pattern: str, root: str = "", max_results: int = 50) -> str:
    search_root = (ws / root).resolve() if root else ws.resolve()
    if not str(search_root).startswith(str(ws.resolve())):
        return "Access denied for this path"
    if not search_root.exists():
        return f"Directory not found: {root or '/'}"

    max_results = max(1, min(int(max_results), 200))
    matches: list[str] = []

    if shutil.which("rg"):
        try:
            proc = subprocess.run(
                [
                    "rg",
                    "--line-number",
                    "--color",
                    "never",
                    "--max-count",
                    str(max_results),
                    pattern,
                    str(search_root),
                ],
                capture_output=True,
                text=True,
                check=False,
            )
            if proc.stdout.strip():
                for line in proc.stdout.splitlines():
                    normalized = line.replace(str(ws.resolve()) + os.sep, "")
                    matches.append(normalized)
            elif proc.returncode not in (0, 1):
                return f"Grep search failed: {proc.stderr.strip()[:200]}"
        except Exception as e:
            return f"Grep search failed: {e}"
    else:
        try:
            for path in sorted(search_root.rglob("*")):
                if len(matches) >= max_results:
                    break
                if not path.is_file():
                    continue
                try:
                    with path.open("r", encoding="utf-8", errors="replace") as handle:
                        for idx, line in enumerate(handle, start=1):
                            if pattern in line:
                                matches.append(f"{path.relative_to(ws).as_posix()}:{idx}:{line.strip()}")
                                if len(matches) >= max_results:
                                    break
                except Exception as _read_err:
                    logger.debug("[Workspace] grep: skipped file %s: %s", path, _read_err)
                    continue
        except Exception as e:
            return f"Grep search failed: {e}"

    if not matches:
        return f"🔎 No matches for '{pattern}'"
    lines = [f"🔎 Grep results for '{pattern}' ({len(matches)} match(es)):"]
    lines.extend(f"- {match}" for match in matches[:max_results])
    return "\n".join(lines)


def _tool_search(ws: Path, query: str = "") -> str:
    packs = iter_tool_packs(query)
    registry = _build_skill_registry(ws)
    normalized = query.strip().lower()
    matching_skills = [
        skill
        for skill in (registry.resolve(name) for name in registry.names())
        if not normalized
        or normalized in skill.metadata.name.lower()
        or normalized in skill.metadata.description.lower()
        or any(normalized in tool.lower() for tool in skill.metadata.declared_tools)
    ]

    lines = [
        "Tool search only returns delayed capability summaries. It does not auto-load tools.",
    ]
    if packs:
        lines.append("")
        lines.append("Available packs:")
        for pack in packs:
            tools = ", ".join(pack.tools)
            lines.append(
                f"- {pack.name}: {pack.summary} | tools: {tools} | activation: {pack.activation_mode}"
            )
    if matching_skills:
        lines.append("")
        lines.append("Matching skills:")
        for skill in matching_skills[:20]:
            declared = ", ".join(skill.metadata.declared_tools) if skill.metadata.declared_tools else "no declared tools"
            lines.append(
                f"- {skill.metadata.name}: {skill.metadata.description} | declared tools: {declared}"
            )
    if len(lines) == 1:
        return f"🔎 No delayed tools or skills matched '{query}'"
    return "\n".join(lines)


def _delete_file(ws: Path, rel_path: str) -> str:
    protected = {"tasks.json", "soul.md"}
    if rel_path.strip("/") in protected:
        return f"{rel_path} cannot be deleted (protected)"

    file_path = (ws / rel_path).resolve()
    if not str(file_path).startswith(str(ws.resolve())):
        return "Access denied for this path"
    if not file_path.exists():
        return f"File not found: {rel_path}"

    try:
        if file_path.is_dir():
            shutil.rmtree(file_path)
            return f"✅ Deleted directory {rel_path}"
        file_path.unlink()
        return f"✅ Deleted {rel_path}"
    except Exception as e:
        return f"Delete failed: {e}"
